import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

base = Path(r"C:\Users\Ale\Documents\Zoom\2026-02-20 10.09.18 Reunión de Zoom de Alejandro Rebolledo\analysis")
steps = json.loads((base/"important_steps_detailed.json").read_text(encoding="utf-8"))
out = base/"Guia_registro_sensor_suelo_Wincontop_v2_detallada.pptx"

prs = Presentation()

# portada
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "Guía: Registro sensor de suelo Wincontop"
s.placeholders[1].text = "Versión detallada con evidencia visual por timestamp"

# checklist
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Checklist operativo"
tf = s.placeholders[1].text_frame
tf.clear()
items = [
"1) Entrar al administrador e identificar proyecto/dispositivo base",
"2) Validar tipo de sensor y asociaciones existentes",
"3) Crear/ajustar relación en Variables en sensores",
"4) Confirmar unidad de CE (µS/cm)",
"5) Duplicar/configurar dispositivo y segundo sensor",
"6) Verificar orden de variables en link",
"7) Guardar y validar telemetría"
]
for i,it in enumerate(items):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.text=it

# slides por paso con pantallazo
for r in steps:
    sl = prs.slides.add_slide(prs.slide_layouts[5])
    sl.shapes.title.text = f"Paso {r['step']} - {r['title']}"

    # caja de acción y timestamp
    tx = sl.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.2), Inches(1.2))
    tf = tx.text_frame
    mm = int(r['timestamp_s']//60); ss=int(r['timestamp_s']%60)
    tf.text = f"Timestamp: {mm:02d}:{ss:02d} | Acción: {r['action']}"
    tf.paragraphs[0].font.size = Pt(16)

    # evidencia transcripción
    tx2 = sl.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.2), Inches(1.0))
    tf2 = tx2.text_frame
    tf2.text = f"Evidencia: {r['evidence_text']}"
    tf2.paragraphs[0].font.size = Pt(12)

    # imagen por timestamp
    sec = int(r['timestamp_s'])
    frame = base/"frames"/f"step_{r['step']:02d}_{sec:04d}s.jpg"
    if frame.exists():
        sl.shapes.add_picture(str(frame), Inches(1.0), Inches(3.0), height=Inches(4.0))

# cierre
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Validación final"
s.placeholders[1].text = "• Confirmar CE, temperatura y humedad de suelo\n• Confirmar mapeo y orden del link\n• Confirmar recepción de datos en histórico"

prs.save(out)
print(out)
