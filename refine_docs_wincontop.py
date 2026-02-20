import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

base = Path(r"C:\Users\Ale\Documents\Zoom\2026-02-20 10.09.18 Reunión de Zoom de Alejandro Rebolledo\analysis")
frames_dir = base / "frames"
frames_dir.mkdir(exist_ok=True)
video = Path(r"C:\Users\Ale\Documents\Zoom\2026-02-20 10.09.18 Reunión de Zoom de Alejandro Rebolledo\video1675864608.mp4")
steps_json = base / "important_steps.json"
transcript_json = base / "transcript.json"

j = json.loads(transcript_json.read_text(encoding='utf-8'))
segments = j.get('segments',[])

def nearest_text(t):
    cand = min(segments, key=lambda s: abs(float(s.get('start',0))-t))
    return cand.get('text','').strip(), float(cand.get('start',0)), float(cand.get('end',0))

# Curated timeline from meeting topics (manually distilled from transcript)
plan = [
(21, "Inicio del flujo", "Entrar al administrador para registrar sensor de suelo y revisar estado actual."),
(46, "Sensor tipo vs sensor", "Aclarar que primero se valida/crea el tipo de sensor y luego el sensor/dispositivo."),
(72, "Reutilizar entidades existentes", "Si el sensor ya existe, no recrearlo; completar asociaciones faltantes."),
(193, "Diagnóstico principal", "La conductividad no aparecía por asociación incompleta de variables."),
(232, "Unidad de medida", "Conductividad en µS/cm según datasheet."),
(264, "Variable existente", "Revisar variable existente de conductividad y ajustar etiqueta 'ambiental' si aplica."),
(324, "Relación variable-sensor", "Ir a 'Variables en sensores' y crear relación nueva para sensor de suelo Wincontop."),
(404, "Dispositivo objetivo", "Usar/duplicar dispositivo base y agregar segundo sensor (caso Soil 9)."),
(497, "Duplicación práctica", "Duplicar acelera, luego limpiar variables no usadas (p. ej. módem/GPS)."),
(579, "Composición final", "Confirmar sensores: DHT22 + suelo + batería + módem + segundo sensor."),
(648, "Mapeo en enlace", "Validar orden de variables en el link: temperatura aire, CE suelo, temperatura suelo, humedad suelo, etc."),
(703, "Cierre operativo", "Guardar, duplicar estaciones necesarias y validar telemetría/pop-up de creación."),
]

# Save improved steps json
rows=[]
for i,(t,title,desc) in enumerate(plan, start=1):
    txt,st,en = nearest_text(t)
    rows.append({"step":i,"timestamp_s":t,"title":title,"action":desc,"evidence_text":txt,"seg_start":st,"seg_end":en})

(base/"important_steps_detailed.json").write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding='utf-8')

# Markdown SOP
md = []
md.append("# SOP - Registro de sensor de suelo Wincontop\n")
md.append("Documento corregido y ampliado a partir de la reunión de Zoom.\n")
md.append("## Alcance\n")
md.append("Configuración de sensor de suelo Wincontop con variables de humedad, temperatura y conductividad, incluyendo asociación en plataforma, mapeo de enlace y validación final.\n")
md.append("## Pre-requisitos\n")
md.append("- Acceso al administrador de la plataforma\n- Proyecto de monitoreo de suelos creado\n- Dispositivo base existente (recomendado para duplicar)\n- Datasheet de Wincontop para confirmar unidades\n")
md.append("## Checklist detallado\n")
for r in rows:
    mm = int(r['timestamp_s']//60)
    ss = int(r['timestamp_s']%60)
    md.append(f"### {r['step']}. {r['title']} ({mm:02d}:{ss:02d})\n")
    md.append(f"- **Acción:** {r['action']}\n")
    md.append(f"- **Evidencia (transcripción):** \"{r['evidence_text']}\"\n")

md.append("## Reglas clave levantadas de la reunión\n")
md.append("- Si el sensor ya existe, no recrearlo: completar asociaciones pendientes.\n")
md.append("- La falta de conductividad suele ser problema de asociación/mapeo, no necesariamente del hardware.\n")
md.append("- Unidad de conductividad: **µS/cm**.\n")
md.append("- En duplicación de dispositivos, revisar variables heredadas y limpiar las no usadas.\n")
md.append("- El orden del link debe coincidir con el orden de la tabla/variables en el sistema.\n")

md.append("## Validación final\n")
md.append("1. Confirmar recepción de temperatura de aire (DHT22).\n2. Confirmar CE, temperatura y humedad de suelo.\n3. Confirmar variables auxiliares (batería/módem) según necesidad.\n4. Verificar que la estación duplicada tenga ID/link actualizado.\n")

(base/"GUIA_PASO_A_PASO_WINCONTOP.md").write_text("\n".join(md), encoding='utf-8')

# Update README with fuller content
readme = f'''# Guía de registro de sensor de suelo Wincontop (material para GitHub)

Este directorio contiene entregables corregidos y ampliados de la reunión de Zoom.

## Contenido

- `transcript.txt`  
  Transcripción completa en texto.
- `transcript.json`  
  Transcripción con segmentos y timestamps.
- `important_steps.json`  
  Pasos relevantes detectados automáticamente (versión breve).
- `important_steps_detailed.json`  
  Pasos ampliados y curados (versión recomendada).
- `GUIA_PASO_A_PASO_WINCONTOP.md`  
  SOP detallado para operación en terreno.
- `frames/`  
  Capturas de pantalla del video en momentos clave.
- `Guia_registro_sensor_suelo_Wincontop.pptx`  
  Presentación final breve.
- `Guia_registro_sensor_suelo_Wincontop_v2_detallada.pptx`  
  Presentación corregida y ampliada.

## Objetivo

Documentar un flujo operativo completo para configurar sensor de suelo Wincontop en plataforma: tipo de sensor, asociación de variables, unidad de CE (µS/cm), mapeo de link, duplicación de dispositivo y validación final.

## Estructura sugerida del repositorio

```text
/docs/wincontop-sensor-suelo/
  README.md
  GUIA_PASO_A_PASO_WINCONTOP.md
  transcript.txt
  transcript.json
  important_steps.json
  important_steps_detailed.json
  Guia_registro_sensor_suelo_Wincontop.pptx
  Guia_registro_sensor_suelo_Wincontop_v2_detallada.pptx
  /frames/
```

## Recomendaciones antes de publicar

- Corregir frases de ASR que puedan estar ambiguas (audio con ruido).
- Revisar y anonimizar nombres o datos sensibles.
- Confirmar que cada captura coincida con el paso descrito.

## Estado

- Whisper local: configurado en `C:\\whisper`
- Transcripción: completada
- Capturas clave: completadas
- Documentación detallada (MD): completada
- PPT v2 ampliada: completada
'''
(base/"README_GITHUB.md").write_text(readme, encoding='utf-8')

# Build detailed PPT v2
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Guía corregida: sensor de suelo Wincontop"
slide.placeholders[1].text = "Reunión Zoom -> versión detallada para ejecución y documentación"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Qué se corrigió"
slide.placeholders[1].text = "• Se eliminó referencia a TEROS\n• Se amplió el flujo completo\n• Se agregó mapeo de variables/link\n• Se agregó checklist de validación final"

for r in rows:
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    sl.shapes.title.text = f"Paso {r['step']} - {r['title']}"
    tf = sl.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Acción: {r['action']}"
    p.font.size = Pt(20)
    p2 = tf.add_paragraph()
    p2.text = f"Evidencia: {r['evidence_text']}"
    p2.level = 1
    p2.font.size = Pt(14)
    p3 = tf.add_paragraph()
    m = int(r['timestamp_s']//60); s=int(r['timestamp_s']%60)
    p3.text = f"Timestamp aproximado: {m:02d}:{s:02d}"
    p3.level = 1
    p3.font.size = Pt(14)

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Checklist de cierre"
slide.placeholders[1].text = "1) CE en µS/cm\n2) Variables asociadas correctamente\n3) Orden del link validado\n4) Estación/dispositivo operativo\n5) Datos visibles en histórico"

out = base / "Guia_registro_sensor_suelo_Wincontop_v2_detallada.pptx"
prs.save(out)
print("done", out)
