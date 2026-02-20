import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

base = Path(r"C:\Users\Ale\Documents\Zoom\2026-02-20 10.09.18 Reunión de Zoom de Alejandro Rebolledo\analysis")
out = base / "Guia_registro_sensor_suelo_Wincontop.pptx"
steps = json.loads((base/"important_steps.json").read_text(encoding="utf-8"))

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Guía de registro de sensor de suelo Wincontop"
slide.placeholders[1].text = "Basado en reunión Zoom (capturas + pasos operativos)"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Checklist operativo"
tf = slide.placeholders[1].text_frame
tf.clear()
items = [
"1) Ingresar al administrador e iniciar registro del sensor.",
"2) Validar/crear tipo de sensor de suelo Wincontop.",
"3) Confirmar variables: humedad de suelo, temperatura y conductividad.",
"4) Ir a Variables en sensores y crear relación nueva.",
"5) Definir unidad de conductividad (µS/cm) según ficha técnica.",
"6) Asociar al dispositivo/estación y guardar.",
"7) Verificar recepción de telemetría y datos históricos."
]
for i,it in enumerate(items):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.text = it

for s in steps:
    sl = prs.slides.add_slide(prs.slide_layouts[5])
    sl.shapes.title.text = f"Paso {s['step']} — {s['timestamp']}"
    tx = sl.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.0), Inches(1.2))
    t = tx.text_frame
    t.text = s.get('text','')
    t.paragraphs[0].font.size = Pt(15)
    fr = Path(s.get('frame',''))
    if fr.exists():
        sl.shapes.add_picture(str(fr), Inches(0.7), Inches(2.0), height=Inches(4.8))

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Nota"
slide.placeholders[1].text = "Corrección aplicada: no corresponde TEROS; el flujo se documenta para sensor de suelo Wincontop."

prs.save(out)
print(out)
