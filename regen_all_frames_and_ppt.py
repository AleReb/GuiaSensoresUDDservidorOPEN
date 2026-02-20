import json, subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

base = Path(r"C:\Users\Ale\Documents\Zoom\2026-02-20 10.09.18 Reunión de Zoom de Alejandro Rebolledo\analysis")
video = Path(r"C:\Users\Ale\Documents\Zoom\2026-02-20 10.09.18 Reunión de Zoom de Alejandro Rebolledo\video1675864608.mp4")
ffmpeg = Path(r"C:\whisper\ffmpeg\ffmpeg-8.0.1-essentials_build\bin\ffmpeg.exe")
steps = json.loads((base/"important_steps_detailed.json").read_text(encoding="utf-8"))
frames = base/"frames"
frames.mkdir(exist_ok=True)

for r in steps:
    step = int(r['step'])
    t = float(r['timestamp_s'])
    hh = int(t//3600); mm = int((t%3600)//60); ss = t%60
    ts = f"{hh:02d}:{mm:02d}:{ss:06.3f}"
    out = frames / f"step_{step:02d}.jpg"
    cmd=[str(ffmpeg),"-y","-ss",ts,"-i",str(video),"-frames:v","1","-q:v","2",str(out)]
    subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=False)

out_ppt = base/"Guia_registro_sensor_suelo_Wincontop_v2_detallada.pptx"
prs = Presentation()

s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "Guía: Registro sensor de suelo Wincontop"
s.placeholders[1].text = "Versión detallada con pantallazo en cada paso"

s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Checklist operativo"
s.placeholders[1].text = "1) Administrador/proyecto\n2) Tipo de sensor\n3) Variables en sensores\n4) Unidad CE µS/cm\n5) Dispositivo/duplicación\n6) Mapeo link\n7) Validación telemetría"

for r in steps:
    sl = prs.slides.add_slide(prs.slide_layouts[5])
    sl.shapes.title.text = f"Paso {r['step']} - {r['title']}"
    tx=sl.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.2), Inches(1.2))
    tf=tx.text_frame
    mm = int(float(r['timestamp_s'])//60); ss=int(float(r['timestamp_s'])%60)
    tf.text = f"Timestamp {mm:02d}:{ss:02d} | {r['action']}"
    tf.paragraphs[0].font.size = Pt(15)
    tx2=sl.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.2), Inches(1.0))
    tf2=tx2.text_frame
    tf2.text = f"Evidencia: {r['evidence_text']}"
    tf2.paragraphs[0].font.size = Pt(12)
    img = frames / f"step_{int(r['step']):02d}.jpg"
    if img.exists() and img.stat().st_size>0:
        sl.shapes.add_picture(str(img), Inches(0.9), Inches(3.0), height=Inches(3.9))

s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Validación final"
s.placeholders[1].text = "Confirmar CE/temperatura/humedad + orden de link + histórico de datos"

prs.save(out_ppt)
print('ok')
