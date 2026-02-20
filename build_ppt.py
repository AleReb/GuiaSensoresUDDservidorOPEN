import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

base = Path(r"C:\Users\Ale\Documents\Zoom\2026-02-20 10.09.18 Reunión de Zoom de Alejandro Rebolledo\analysis")
frames = base / "frames"
out = base / "Guia_registro_sensores_TEROS12.pptx"
steps_json = base / "important_steps.json"

steps = json.loads(steps_json.read_text(encoding="utf-8")) if steps_json.exists() else []

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Guía de registro de sensores (extraída de Zoom)"
slide.placeholders[1].text = "Caso: TEROS 12 / humedad, temperatura y conductividad"

# Process summary slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Checklist operativo"
tf = slide.placeholders[1].text_frame
tf.clear()
checklist = [
"1) Entrar al administrador e iniciar registro del sensor.",
"2) Validar/crear Tipo de sensor (evitar duplicados).",
"3) Confirmar nombre/código y asociación correcta.",
"4) Ir a Variables en sensores.",
"5) Crear relación nueva para conductividad.",
"6) Asignar unidad (µS/cm) según datasheet.",
"7) Asociar al dispositivo/estación (ej. Soil 9).",
"8) Guardar y validar recepción de variables."
]
for i, line in enumerate(checklist):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.level = 0

# Step slides with images
for s in steps:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"Paso {s.get('step')} — {s.get('timestamp')}"

    text = s.get("text", "")
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.0))
    tf = tx.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(16)

    frame = Path(s.get("frame", ""))
    if frame.exists():
        slide.shapes.add_picture(str(frame), Inches(0.7), Inches(2.0), height=Inches(4.8))

# Closing slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Validación final"
slide.placeholders[1].text = "• Revisar que el sensor reporte VWC, temperatura y CE.\n• Confirmar unidad CE en µS/cm.\n• Verificar historial de datos y guardar evidencia."

prs.save(out)
print(out)
